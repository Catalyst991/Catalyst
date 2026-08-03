import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from catalyst.tools.daily_report.excel_reader import Comment
from catalyst.tools.daily_report.report_builder import COLUMN_STYLES, build_report


def make_comment(**overrides):
    defaults = dict(
        date=datetime.date(2026, 7, 27),
        user_name="Design Studio",
        comment="hello world",
        link="https://x.com/1",
        country="KSA",
        follower_count=1296,
        tone="محايد",
    )
    defaults.update(overrides)
    return Comment(**defaults)


def test_builds_a_three_slide_presentation(template_path):
    comments = [make_comment()]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    assert len(prs.slides) == 3


def _find_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    raise AssertionError("no table found on slide")


def test_fills_content_table_and_leaves_unused_rows_blank(template_path):
    comments = [
        make_comment(user_name="Design Studio", comment="hello world"),
        make_comment(user_name="Sultan", comment="second comment"),
    ]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    table = _find_table(prs.slides[1])
    assert [c.text for c in table.rows[1].cells] == [
        "2026-07-27",
        "Design Studio",
        "hello world",
        "https://x.com/1",
        "KSA",
        "1296",
        "محايد",
    ]
    assert table.rows[2].cells[1].text == "Sultan"
    assert table.rows[2].cells[2].text == "second comment"
    for row_index in range(3, len(table.rows)):
        assert all(cell.text == "" for cell in table.rows[row_index].cells)


def test_embedded_newlines_are_normalized_to_a_single_styled_paragraph(template_path):
    comments = [
        make_comment(
            user_name="Design Studio\n",
            comment="سؤال \nلماذا يسمح لأصحاب المحلات\nالإعلان لمنتجاتهم",
        )
    ]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    table = _find_table(prs.slides[1])
    user_name_cell = table.cell(1, 1)
    comment_cell = table.cell(1, 2)

    assert len(user_name_cell.text_frame.paragraphs) == 1
    assert "\n" not in user_name_cell.text
    assert user_name_cell.text == "Design Studio"

    assert len(comment_cell.text_frame.paragraphs) == 1
    assert "\n" not in comment_cell.text
    run = comment_cell.text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(8.0)
    assert run.font.name == "Segoe UI"


def _title_text(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "عنوان 1":
            return shape.text_frame.text
    raise AssertionError("no title shape found")


def test_title_shows_generation_date_not_any_date_in_the_data(template_path):
    comments = [make_comment(date=datetime.date(2020, 1, 1))]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    title = _title_text(prs.slides[0])
    assert "28 يوليو" in title
    assert "يناير" not in title


def test_title_split_across_multiple_runs_is_merged_without_leftovers(template_path):
    prs = Presentation(template_path)
    for shape in prs.slides[0].shapes:
        if shape.name == "عنوان 1":
            paragraph = shape.text_frame.paragraphs[0]
            paragraph.clear()
            paragraph.add_run().text = "تقرير الرصد اليومي – "
            paragraph.add_run().text = "25 مايو"
    prs.save(template_path)

    comments = [make_comment()]
    result = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    title = _title_text(result.slides[0])
    assert title == "تقرير الرصد اليومي – 28 يوليو"


def _title_run(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "عنوان 1":
            return shape.text_frame.paragraphs[0].runs[0]
    raise AssertionError("no title shape found")


def test_title_keeps_its_original_font_and_color(template_path):
    comments = [make_comment()]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    run = _title_run(prs.slides[0])
    assert run.font.name == "DIN Next LT Arabic Medium"
    assert run.font.size == Pt(40)
    assert run.font.color.rgb == RGBColor(0xA8, 0x9F, 0x88)


def _page_number_text(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 3":
            return shape.text_frame.text
    raise AssertionError("no page number shape found")


def test_content_slide_page_number_matches_its_slide_position(template_path):
    comments = [make_comment()]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    assert _page_number_text(prs.slides[1]) == "2"


def test_long_comment_is_truncated_in_the_table(template_path):
    comments = [make_comment(comment="a" * 95)]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    table = _find_table(prs.slides[1])
    assert table.rows[1].cells[2].text == ("a" * 90) + "..."


def test_filled_table_cells_use_the_correct_font_size_and_alignment(template_path):
    comments = [make_comment()]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    table = _find_table(prs.slides[1])
    for col_index, (expected_size, expected_alignment) in enumerate(COLUMN_STYLES):
        cell = table.cell(1, col_index)
        run = cell.text_frame.paragraphs[0].runs[0]
        assert run.font.size == expected_size
        assert run.font.name == "Segoe UI"
        assert cell.text_frame.paragraphs[0].alignment == expected_alignment


def test_all_data_rows_have_the_same_fixed_height(template_path):
    comments = [make_comment(comment="short"), make_comment(comment="a" * 90)]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    table = _find_table(prs.slides[1])
    data_row_heights = [table.rows[i].height for i in range(1, len(table.rows))]
    assert len(set(data_row_heights)) == 1


def test_filled_cells_use_default_wrap_and_autofit_matching_the_template(template_path):
    comments = [make_comment(comment="a" * 90)]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    table = _find_table(prs.slides[1])
    for col_index in range(7):
        text_frame = table.cell(1, col_index).text_frame
        assert text_frame.word_wrap is None
        assert text_frame.auto_size is None


def test_link_cell_has_a_working_hyperlink(template_path):
    comments = [make_comment(link="https://x.com/example/status/123")]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    table = _find_table(prs.slides[1])
    run = table.cell(1, 3).text_frame.paragraphs[0].runs[0]
    assert run.hyperlink.address == "https://x.com/example/status/123"


def test_end_slide_is_unchanged_from_the_template(template_path):
    comments = [make_comment()]

    prs = build_report(template_path, comments, generation_date=datetime.date(2026, 7, 28))

    end_slide = prs.slides[2]
    texts = [s.text_frame.text for s in end_slide.shapes if s.has_text_frame]
    assert "شــكـــــراً" in texts
