from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader

from catalyst.tools.daily_report.pdf_exporter import export_to_pdf


def make_pptx(path, slide_count):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for i in range(slide_count):
        slide = prs.slides.add_slide(blank_layout)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = f"Slide {i + 1}"
    prs.save(path)
    return path


def test_converts_a_pptx_to_a_pdf_with_matching_page_count(tmp_path):
    pptx_path = make_pptx(tmp_path / "report.pptx", slide_count=3)

    pdf_path = export_to_pdf(pptx_path)

    assert pdf_path.exists()
    assert pdf_path.suffix == ".pdf"
    reader = PdfReader(pdf_path)
    assert len(reader.pages) == 3
