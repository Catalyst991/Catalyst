import datetime

import openpyxl
import pytest
from pptx import Presentation
from pypdf import PdfReader

from catalyst.tools.daily_report.excel_reader import MissingSheetError, NoCommentsError
from catalyst.tools.daily_report.pipeline import generate_report

TABLE_HEADERS = ["Date", "User name", "Comment", "Link", "Country", "No. of Followers", "النبرة"]


def make_excel(path, rows):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Users"
    ws.append(TABLE_HEADERS)
    for row in rows:
        ws.append(row)
    wb.save(path)
    return path


def test_generates_a_pptx_file_readable_back_with_correct_contents(tmp_path, template_path):
    excel_path = make_excel(
        tmp_path / "export.xlsx",
        rows=[
            [datetime.date(2026, 7, 27), "Design Studio", "hello world", "https://x.com/1", "KSA", 1296, "محايد"],
        ],
    )
    save_directory = tmp_path / "reports"
    save_directory.mkdir()

    save_paths = generate_report(
        excel_path=excel_path,
        template_path=template_path,
        save_directory=save_directory,
        generation_date=datetime.date(2026, 7, 28),
    )

    assert len(save_paths) == 1
    save_path = save_paths[0]
    assert save_path.exists()
    prs = Presentation(save_path)
    assert len(prs.slides) == 3
    table = next(shape.table for shape in prs.slides[1].shapes if shape.has_table)
    assert table.rows[1].cells[1].text == "Design Studio"


def test_writes_no_file_when_the_excel_file_fails_validation(tmp_path, template_path):
    wb = openpyxl.Workbook()
    wb.active.title = "Official"
    excel_path = tmp_path / "export.xlsx"
    wb.save(excel_path)
    save_directory = tmp_path / "reports"
    save_directory.mkdir()

    with pytest.raises(MissingSheetError):
        generate_report(
            excel_path=excel_path,
            template_path=template_path,
            save_directory=save_directory,
            generation_date=datetime.date(2026, 7, 28),
        )

    assert list(save_directory.iterdir()) == []


def test_writes_no_file_when_the_excel_file_has_zero_comments(tmp_path, template_path):
    excel_path = make_excel(tmp_path / "export.xlsx", rows=[])
    save_directory = tmp_path / "reports"
    save_directory.mkdir()

    with pytest.raises(NoCommentsError):
        generate_report(
            excel_path=excel_path,
            template_path=template_path,
            save_directory=save_directory,
            generation_date=datetime.date(2026, 7, 28),
        )

    assert list(save_directory.iterdir()) == []


def test_saved_filename_matches_the_title_slide_text(tmp_path, template_path):
    excel_path = make_excel(
        tmp_path / "export.xlsx",
        rows=[
            [datetime.date(2026, 7, 27), "Design Studio", "hello world", "https://x.com/1", "KSA", 1296, "محايد"],
        ],
    )
    save_directory = tmp_path / "reports"
    save_directory.mkdir()

    save_paths = generate_report(
        excel_path=excel_path,
        template_path=template_path,
        save_directory=save_directory,
        generation_date=datetime.date(2026, 7, 28),
    )

    assert len(save_paths) == 1
    save_path = save_paths[0]
    assert save_path.parent == save_directory
    assert save_path.suffix == ".pptx"
    assert save_path.stem == "تقرير الرصد اليومي – 28 يوليو"


def test_pdf_only_saves_a_single_pdf_and_leaves_no_pptx_anywhere(tmp_path, template_path):
    excel_path = make_excel(
        tmp_path / "export.xlsx",
        rows=[
            [datetime.date(2026, 7, 27), "Design Studio", "hello world", "https://x.com/1", "KSA", 1296, "محايد"],
        ],
    )
    save_directory = tmp_path / "reports"
    save_directory.mkdir()

    save_paths = generate_report(
        excel_path=excel_path,
        template_path=template_path,
        save_directory=save_directory,
        output_format="pdf",
        generation_date=datetime.date(2026, 7, 28),
    )

    assert len(save_paths) == 1
    pdf_path = save_paths[0]
    assert pdf_path.suffix == ".pdf"
    assert pdf_path.parent == save_directory
    reader = PdfReader(pdf_path)
    assert len(reader.pages) == 3

    report_pptx_files = [p for p in tmp_path.rglob("*.pptx") if p != template_path]
    assert report_pptx_files == []


def test_both_saves_a_pptx_and_a_pdf_with_matching_base_filenames(tmp_path, template_path):
    excel_path = make_excel(
        tmp_path / "export.xlsx",
        rows=[
            [datetime.date(2026, 7, 27), "Design Studio", "hello world", "https://x.com/1", "KSA", 1296, "محايد"],
        ],
    )
    save_directory = tmp_path / "reports"
    save_directory.mkdir()

    save_paths = generate_report(
        excel_path=excel_path,
        template_path=template_path,
        save_directory=save_directory,
        output_format="both",
        generation_date=datetime.date(2026, 7, 28),
    )

    assert len(save_paths) == 2
    suffixes = {p.suffix for p in save_paths}
    assert suffixes == {".pptx", ".pdf"}
    stems = {p.stem for p in save_paths}
    assert stems == {"تقرير الرصد اليومي – 28 يوليو"}
    for path in save_paths:
        assert path.exists()
        assert path.parent == save_directory
