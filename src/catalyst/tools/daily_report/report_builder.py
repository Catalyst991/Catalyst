import copy
from datetime import date

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from catalyst.tools.daily_report.truncate import truncate

COMMENT_MAX_LENGTH = 90
COMMENTS_PER_SLIDE = 7

FONT_NAME = "Segoe UI"

LINK_COLUMN_INDEX = 3

# (font size, alignment, right-to-left) per column, captured from the real template's
# filled cells (identical across every filled row and every content slide in the real
# file), with alignment/rtl adjusted for the Arabic-language columns
COLUMN_STYLES = [
    (Pt(8.19), PP_ALIGN.CENTER, False),  # Date
    (Pt(8.19), PP_ALIGN.CENTER, False),  # User name
    (Pt(8.0), PP_ALIGN.RIGHT, True),  # Comment
    (Pt(7.9), PP_ALIGN.LEFT, False),  # Link
    (Pt(8.19), PP_ALIGN.CENTER, False),  # Country
    (Pt(8.19), PP_ALIGN.CENTER, False),  # No. of Followers
    (Pt(8.0), PP_ALIGN.CENTER, True),  # Tone (النبرة)
]

ARABIC_MONTH_NAMES = {
    1: "يناير",
    2: "فبراير",
    3: "مارس",
    4: "أبريل",
    5: "مايو",
    6: "يونيو",
    7: "يوليو",
    8: "أغسطس",
    9: "سبتمبر",
    10: "أكتوبر",
    11: "نوفمبر",
    12: "ديسمبر",
}


def _format_report_date(generation_date: date) -> str:
    return f"{generation_date.day} {ARABIC_MONTH_NAMES[generation_date.month]}"


def _set_title_date(slide, generation_date: date) -> None:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "عنوان 1":
            runs = shape.text_frame.paragraphs[0].runs
            full_text = "".join(run.text for run in runs)
            prefix = full_text.split("–")[0].strip()

            target_index = next((i for i, run in enumerate(runs) if "–" in run.text), len(runs) - 1)
            for i, run in enumerate(runs):
                run.text = f"{prefix} – {_format_report_date(generation_date)}" if i == target_index else ""
            return
    raise ValueError("no title shape found on start slide")


def _normalize_whitespace(text: str) -> str:
    return " ".join(text.split())


def _find_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    raise ValueError("no table found on slide")


def _set_cell_text(cell, text: str, size: Pt, alignment: PP_ALIGN, rtl: bool) -> None:
    cell.text = text
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = alignment
    if rtl:
        paragraph._pPr.set("rtl", "1")
    run = paragraph.runs[0]
    run.font.size = size
    run.font.name = FONT_NAME


def _fill_content_table(slide, comments):
    table = _find_table(slide)

    data_row_height = table.rows[1].height
    for row_index in range(1, len(table.rows)):
        table.rows[row_index].height = data_row_height

    for row_index, comment in enumerate(comments, start=1):
        values = [
            comment.date.isoformat(),
            _normalize_whitespace(comment.user_name),
            truncate(_normalize_whitespace(comment.comment), COMMENT_MAX_LENGTH),
            _normalize_whitespace(comment.link),
            _normalize_whitespace(comment.country),
            str(comment.follower_count),
            _normalize_whitespace(comment.tone),
        ]
        for col_index, value in enumerate(values):
            size, alignment, rtl = COLUMN_STYLES[col_index]
            cell = table.cell(row_index, col_index)
            _set_cell_text(cell, value, size, alignment, rtl)
            if col_index == LINK_COLUMN_INDEX:
                cell.text_frame.paragraphs[0].runs[0].hyperlink.address = value


def _chunk(items: list, size: int) -> list[list]:
    chunks = [items[i : i + size] for i in range(0, len(items), size)]
    return chunks or [[]]


def _duplicate_slide(prs: Presentation, slide):
    duplicate = prs.slides.add_slide(slide.slide_layout)
    for shape in list(duplicate.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in slide.shapes:
        duplicate.shapes._spTree.append(copy.deepcopy(shape._element))
    return duplicate


def _move_slide_to_end(prs: Presentation, slide) -> None:
    slide_id_list = prs.slides._sldIdLst
    element = next(el for el in slide_id_list if el.get("id") == str(slide.slide_id))
    slide_id_list.remove(element)
    slide_id_list.append(element)


def _build_content_slides(prs: Presentation, comments: list) -> None:
    end_slide = prs.slides[2]
    content_template = prs.slides[1]
    chunks = _chunk(comments, COMMENTS_PER_SLIDE)

    content_slides = [content_template]
    for _ in chunks[1:]:
        content_slides.append(_duplicate_slide(prs, content_template))

    _move_slide_to_end(prs, end_slide)

    for slide, chunk in zip(content_slides, chunks):
        _fill_content_table(slide, chunk)


def _set_page_numbers(prs: Presentation) -> None:
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == "TextBox 3":
                shape.text_frame.text = str(slide_index)


def build_report(template_path, comments, generation_date: date) -> Presentation:
    prs = Presentation(template_path)
    _set_title_date(prs.slides[0], generation_date)
    _build_content_slides(prs, comments)
    _set_page_numbers(prs)
    return prs
