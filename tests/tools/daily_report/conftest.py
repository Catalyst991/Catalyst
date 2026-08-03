import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

TABLE_HEADERS = ["Date", "User name", "Comment", "Link", "Country", "No. of Followers", "النبرة"]


def make_template(path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    start_slide = prs.slides.add_slide(blank_layout)
    title_box = start_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_box.name = "عنوان 1"
    title_box.text_frame.text = "تقرير الرصد اليومي – 25 مايو"
    title_run = title_box.text_frame.paragraphs[0].runs[0]
    title_run.font.name = "DIN Next LT Arabic Medium"
    title_run.font.size = Pt(40)
    title_run.font.color.rgb = RGBColor(0xA8, 0x9F, 0x88)

    content_slide = prs.slides.add_slide(blank_layout)
    page_num_box = content_slide.shapes.add_textbox(Inches(9), Inches(0.2), Inches(0.5), Inches(0.3))
    page_num_box.name = "TextBox 3"
    page_num_box.text_frame.text = "X"

    table_shape = content_slide.shapes.add_table(8, 7, Inches(0.5), Inches(1), Inches(9), Inches(5))
    table_shape.name = "Table 4"
    table = table_shape.table
    for col, header in enumerate(TABLE_HEADERS):
        table.cell(0, col).text = header

    end_slide = prs.slides.add_slide(blank_layout)
    thanks_box = end_slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(1))
    thanks_box.text_frame.text = "شــكـــــراً"

    prs.save(path)
    return path


@pytest.fixture
def template_path(tmp_path):
    return make_template(tmp_path / "template.pptx")
